import os
import glob
import requests
import mimetypes
from datetime import datetime
import pypdf 
import docx
from pptx import Presentation
from fastapi import UploadFile

from langchain_openai import ChatOpenAI, OpenAIEmbeddings
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_core.documents import Document
from langchain_chroma import Chroma
from typing import Iterator, List
from langchain_text_splitters import RecursiveCharacterTextSplitter

from langchain_core.vectorstores import InMemoryVectorStore

from .logging_config import get_logger

# Initialize logger for this module
logger = get_logger(__name__)

class LocalFileService:
    def __init__(self, root_dir: str, db_manager=None):
        self.root_dir = root_dir
        self.db_manager = db_manager
        logger.info(f"Initializing LocalFileService with root_dir: {root_dir}")
        # Ensure directory exists for PoC
        if not os.path.exists(self.root_dir):
            logger.info(f"Creating root directory: {root_dir}")
            os.makedirs(self.root_dir)
        logger.info("LocalFileService initialized successfully")

    def list_files(self):
        logger.debug("Listing files in data directory")
        files_data = []
        # Recursive search for relevant extensions
        extensions = ['*.md', '*.pdf', '*.docx', '*.pptx', '*.xlsx', '*.txt']
        
        # Helper to check summaries efficiently
        summaries_set = set()
        if self.db_manager:
            summaries_set = set(self.db_manager.get_files_with_summaries())

        # globe recursive needs python 3.10+ for 'root_dir/**/ext' or manual walk
        # Using os.walk for better compatibility and control
        for root, dirs, files in os.walk(self.root_dir):
            for file in files:
                if self._match_ext(file, extensions):
                    full_path = os.path.join(root, file)
                    rel_path = os.path.relpath(full_path, self.root_dir)
                    
                    stat = os.stat(full_path)
                    last_modified_dt = datetime.fromtimestamp(stat.st_mtime)
                    file_type = os.path.splitext(file)[1][1:]
                    
                    file_info = {
                        "name": file,
                        "path": rel_path,
                        "type": file_type,
                        "size": stat.st_size,
                        "modified": last_modified_dt.strftime('%Y-%m-%d %H:%M:%S'),
                        "has_summary": rel_path in summaries_set
                    }
                    files_data.append(file_info)
                    
                    # Sync to DB if available
                    if self.db_manager:
                        self.db_manager.upsert_file_metadata(
                            path=rel_path,
                            filename=file,
                            last_modified=last_modified_dt,
                            size=stat.st_size,
                            file_type=file_type
                        )
                        file_info["tags"] = self.db_manager.get_file_tags(rel_path)
        logger.debug(f"Found {len(files_data)} files")
        return files_data

    def _match_ext(self, filename, extensions):
        # Simplified matcher
        ext = os.path.splitext(filename)[1]
        return any(pat.endswith(ext) for pat in extensions)

    def read_pdf(self, file_path):
        """Read content from a PDF file.
        
        Args:
            file_path: Can be a string path or a file-like object
        
        Returns:
            str: Extracted text content from the PDF
        """
        try:
            text = ""
            reader = pypdf.PdfReader(file_path)
            for page in reader.pages:
                text += page.extract_text() + "\n"
            logger.debug(f"Successfully read PDF, {len(reader.pages)} pages")
            return text
        except Exception as e:
            logger.error(f"Failed to read PDF: {str(e)}")
            raise Exception(f"Failed to read PDF: {e}")

    def read_docx(self, file_path):
        """Read content from a DOCX file.
        
        Args:
            file_path: Can be a string path or a file-like object
        
        Returns:
            str: Extracted text content from the DOCX
        """
        try:
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            logger.debug(f"Successfully read DOCX, {len(doc.paragraphs)} paragraphs")
            return text
        except Exception as e:
            logger.error(f"Failed to read DOCX: {str(e)}")
            raise Exception(f"Failed to read DOCX: {e}")

    def read_pptx(self, file_path):
        """Read content from a PPTX file.
        
        Args:
            file_path: Can be a string path or a file-like object
        
        Returns:
            str: Extracted text content from the PPTX
        """
        try:
            prs = Presentation(file_path)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            logger.debug(f"Successfully read PPTX, {len(prs.slides)} slides")
            return "\n".join(text)
        except Exception as e:
            logger.error(f"Failed to read PPTX: {str(e)}")
            raise Exception(f"Failed to read PPTX: {e}")

    def get_content(self, rel_path):
        logger.debug(f"Getting content for file: {rel_path}")
        full_path = os.path.join(self.root_dir, rel_path)
        if not os.path.exists(full_path):
            logger.warning(f"File not found: {rel_path}")
            raise FileNotFoundError(f"File not found: {rel_path}")

        ext = os.path.splitext(full_path)[1].lower()
        logger.debug(f"File extension: {ext}")
        
        if ext in ['.md', '.txt']:
            try:
                with open(full_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                logger.debug(f"Successfully read text file: {rel_path}, size: {len(content)} chars")
                return {"content": content, "type": "text"}
            except Exception as e:
                logger.error(f"Failed to read text file {rel_path}: {str(e)}")
                raise Exception(f"Failed to read text file: {e}")
        
        elif ext == '.pdf':
            text = self.read_pdf(full_path)
            return {"content": text, "type": "text"}

        elif ext == '.docx':
            text = self.read_docx(full_path)
            return {"content": text, "type": "text"}

        elif ext == '.pptx':
            text = self.read_pptx(full_path)
            return {"content": text, "type": "text"}
                
        else:
            # For binary files not yet supported (e.g. xlsx), return message
            logger.debug(f"Binary file not displayable: {rel_path}")
            raise Exception("Binary file content not displayable in text view yet.")
    
    def get_uploaded_file_content(self, file: UploadFile):
        """Get content from an uploaded file object.
        
        Args:
            file: File-like object with filename and read() method
        
        Returns:
            dict: Dictionary with 'content' and 'type' keys
        """
        logger.info(f"Getting content from uploaded file: {file.filename}")
        ext = os.path.splitext(file.filename)[1].lower()
        logger.debug(f"File extension: {ext}")
        
        try:
            if ext in ['.md', '.txt']:
                content = file.file.read()
                if isinstance(content, bytes):
                    content = content.decode('utf-8', errors='ignore')
                return {"content": content, "type": "text"}
            elif ext == '.pdf':
                content = self.read_pdf(file.file)
                return {"content": content, "type": "text"}
            elif ext == '.docx':
                content = self.read_docx(file.file)
                return {"content": content, "type": "text"}
            elif ext == '.pptx':
                content = self.read_pptx(file.file)
                return {"content": content, "type": "text"}
            else:
                logger.debug(f"Binary file not displayable: {file.filename}")
                raise Exception("Binary file content not displayable in text view yet.")
        except Exception as e:
            logger.error(f"Error getting content from uploaded file: {str(e)}", exc_info=True)
            raise e

    def save_content(self, filename: str, content: str):
        # Basic sanitization could be done here
        safe_name = os.path.basename(filename)
        full_path = os.path.join(self.root_dir, safe_name)
        logger.info(f"Saving content to file: {safe_name}")
        
        try:
            with open(full_path, 'w', encoding='utf-8') as f:
                f.write(content)
            
            # Update DB
            if self.db_manager:
                 stat = os.stat(full_path)
                 self.db_manager.upsert_file_metadata(
                    path=safe_name, # since we save to root of data dir
                    filename=safe_name,
                    last_modified=datetime.now(),
                    size=stat.st_size,
                    file_type=os.path.splitext(safe_name)[1][1:]
                )

            logger.info(f"Successfully saved file: {safe_name}, size: {len(content)} chars")
            return {"success": True, "path": safe_name}
        except Exception as e:
            logger.error(f"Failed to save file {safe_name}: {str(e)}")
            return {"error": str(e)}

    def save_upload(self, file_obj, filename: str):
        """Save an uploaded binary file stream."""
        safe_name = os.path.basename(filename)
        full_path = os.path.join(self.root_dir, safe_name)
        logger.info(f"Saving uploaded file: {safe_name}")
        
        try:
            bytes_written = 0
            with open(full_path, 'wb') as f:
                while contents := file_obj.read(1024 * 1024): # Read in chunks
                    f.write(contents)
                    bytes_written += len(contents)
            
            # Update DB
            if self.db_manager:
                 stat = os.stat(full_path)
                 self.db_manager.upsert_file_metadata(
                    path=safe_name, 
                    filename=safe_name,
                    last_modified=datetime.now(),
                    size=stat.st_size,
                    file_type=os.path.splitext(safe_name)[1][1:]
                )
            
            logger.info(f"Successfully saved uploaded file: {safe_name}, size: {bytes_written} bytes")
            return {"success": True, "path": safe_name}
        except Exception as e:
            logger.error(f"Failed to save uploaded file {safe_name}: {str(e)}")
            return {"error": str(e)}

    def delete_file(self, rel_path: str):
        """Delete a file from disk and DB."""
        logger.info(f"Deleting file: {rel_path}")
        full_path = os.path.join(self.root_dir, rel_path)
        
        # Security check: prevent directory traversal
        if ".." in rel_path or not os.path.abspath(full_path).startswith(os.path.abspath(self.root_dir)):
             logger.warning(f"Access denied for file deletion (directory traversal): {rel_path}")
             return {"error": "Access denied"}
        
        try:
            if os.path.exists(full_path):
                os.remove(full_path)
                
                # Update DB
                if self.db_manager:
                    self.db_manager.delete_file(rel_path)
                
                logger.info(f"Successfully deleted file: {rel_path}")
                return {"success": True, "message": f"Deleted {rel_path}"}
            else:
                logger.warning(f"File not found for deletion: {rel_path}")
                return {"error": "File not found"}
        except Exception as e:
            logger.error(f"Failed to delete file {rel_path}: {str(e)}")
            return {"error": f"Failed to delete file: {str(e)}"}

class HedgeDocService:
    def fetch_content(self, url: str, cookie: str = None):
        logger.info(f"Fetching HedgeDoc content from URL: {url}")
        # HedgeDoc usually exposes raw content at /download or if it's already a raw link
        # If user provides view url like https://demo.hedgedoc.org/s/By-Q ...
        # We might need to transform it.
        # For now, assume user provides the link and we try to append /download if not present
        
        target_url = url
        headers = {}
        if cookie:
            headers['Cookie'] = cookie
            logger.debug("Using authentication cookie for HedgeDoc request")

        if not url.endswith('/download') and not '/raw' in url:
             # Basic heuristic: append /download
             # Note: This depends on the specific HedgeDoc instance.
             # Standard HedgeDoc: https://md.example.com/s/Features -> https://md.example.com/Features/download
             # Or https://md.example.com/Features -> https://md.example.com/Features/download
             pass
              
        # Best effort: try adding /download if just fetching the url returns html
        try:
            resp = requests.get(url, headers=headers)
            if resp.status_code != 200:
                # If forbidden, maybe it's private and we have a cookie, but url is wrong 
                # or cookie is invalid.
                logger.warning(f"HedgeDoc request failed with status {resp.status_code}")
                return None
            
            # If Content-Type is text/markdown or text/plain, good.
            ct = resp.headers.get('Content-Type', '')
            if 'text/html' in ct:
                 logger.debug("Received HTML, trying /download endpoint")
                 # Try appending /download
                 download_url = url.rstrip('/') + '/download'
                 resp2 = requests.get(download_url, headers=headers)
                 if resp2.status_code == 200:
                     logger.info(f"Successfully fetched HedgeDoc content from download URL")
                     return resp2.text
                 logger.warning("Could not extract markdown from HedgeDoc")
                 return f"Could not extract markdown. Retrieved HTML from {url}"
            
            logger.info(f"Successfully fetched HedgeDoc content, size: {len(resp.text)} chars")
            return resp.text
        except Exception as e:
            logger.error(f"Error fetching HedgeDoc content: {str(e)}", exc_info=True)
            return f"Error fetching HedgeDoc: {str(e)}"

    def fetch_history(self, base_url: str, cookie: str):
        """
        Fetch history from /history endpoint.
        Requires authenticated session cookie.
        """
        logger.info(f"Fetching HedgeDoc history from: {base_url}")
        if not base_url:
            logger.warning("No base URL provided for HedgeDoc history")
            return None
            
        history_url = base_url.rstrip('/') + '/history'
        headers = {
            'Cookie': cookie,
            'User-Agent': 'Mozilla/5.0' # sometimes needed
        }
        
        try:
            resp = requests.get(history_url, headers=headers)
            if resp.status_code == 200:
                try:
                    data = resp.json()
                    # HedgeDoc /history returns a list of notes
                    # Each note has 'id', 'text' (title), 'time' (last visited/edit)
                    history = data.get('history', [])
                    logger.info(f"Successfully fetched HedgeDoc history, {len(history)} items")
                    return history
                except:
                     logger.error("Failed to parse HedgeDoc history JSON")
                     return {"error": "Failed to parse history JSON", "raw": resp.text[:200]}
            elif resp.status_code == 403:
                logger.warning("HedgeDoc history request forbidden (invalid cookie)")
                return {"error": "Forbidden. Cookie might be invalid."}
            else:
                logger.warning(f"HedgeDoc history request failed with status {resp.status_code}")
                return {"error": f"Failed to fetch history: {resp.status_code}"}
        except Exception as e:
            logger.error(f"Error fetching HedgeDoc history: {str(e)}", exc_info=True)
            return {"error": str(e)}

class GitHubService:
    def get_user_events(self, username: str):
        logger.info(f"Fetching GitHub events for user: {username}")
        url = f"https://api.github.com/users/{username}/events/public"
        try:
            resp = requests.get(url)
            if resp.status_code == 200:
                events = resp.json()
                logger.info(f"Successfully fetched {len(events)} GitHub events for user: {username}")
                # Summarize standard push events
                summary = []
                for e in events[:10]: # Last 10 events
                    etype = e.get('type')
                    repo = e.get('repo', {}).get('name')
                    created_at = e.get('created_at')
                    if etype == 'PushEvent':
                        commits = len(e.get('payload', {}).get('commits', []))
                        summary.append(f"Pushed {commits} commits to {repo} at {created_at}")
                    elif etype == 'CreateEvent':
                         ref_type = e.get('payload', {}).get('ref_type')
                         summary.append(f"Created {ref_type} in {repo} at {created_at}")
                    else:
                        summary.append(f"{etype} at {repo} at {created_at}")
                return {"events": summary, "raw_count": len(events)}
            logger.warning(f"GitHub API returned status {resp.status_code} for user: {username}")
            return None
        except Exception as e:
            logger.error(f"Error fetching GitHub events for user {username}: {str(e)}", exc_info=True)
            return {"error": str(e)}

class LLMService:
    def __init__(self, base_url: str = "http://host.docker.internal:1234/v1", db_manager=None, local_file_service=None):
        self.base_url = base_url
        self.db_manager = db_manager
        self.local_file_service = local_file_service
        logger.info(f"Initializing LLMService with base_url: {base_url}")
        self.llm = ChatOpenAI(
            base_url=self.base_url,
            api_key="lm-studio",
            model="local-model", # Uses default model from LM Studio at this time. TODO: make configurable
            temperature=0.7
        )
        logger.info("LLMService initialized successfully")
    
    def get_models(self):
        logger.debug(f"Fetching models from LLM service: {self.base_url}")
        resp = requests.get(f"{self.base_url}/models")
        if resp.status_code != 200:
            logger.error(f"Failed to get models: {resp.text}")
            raise Exception(f"Failed to get models: {resp.text}")
        models = resp.json()['data']
        logger.debug(f"Found {len(models)} models")
        return models
    

    def get_embedding_models(self):
        models = self.get_models()
        return [m for m in models if "embed" in m['id']]

    def generate_summary(self, content: str) -> str:
        """Generates a summary for the given text content."""
        logger.debug(f"Generating summary for content of length: {len(content)}")
        try:
            messages = [
                SystemMessage(content="You are a helpful assistant that summarizes documents efficiently."),
                HumanMessage(content=f"Please provide a concise summary of the following document:\n\n{content[:8000]}") # Truncate to avoid context window issues if huge
            ]
            response = self.llm.invoke(messages)
            logger.info("Summary generated successfully")
            return response.content
        except Exception as e:
            logger.error(f"Error generating summary: {str(e)}", exc_info=True)
            return f"Error generating summary: {str(e)}"

    def process_file(self, path: str):
        """Reads file, generates summary, and stores in DB."""
        if not self.local_file_service or not self.db_manager:
            return {"error": "Services not fully initialized"}

        # 1. Get Content
        file_data = self.local_file_service.get_content(path)
        if "error" in file_data:
            return file_data
        
        if file_data.get("type") != "text":
             return {"error": "Only text files can be summarized currently."}
            
        content = file_data.get("content")
        if not content:
             return {"error": "File is empty."}

        # 2. Generate Summary
        summary = self.generate_summary(content)
        
        # 3. Save to DB
        # Tags could be extracted by LLM too, but for now we leave empty
        self.db_manager.save_summary(
            path=path,
            summary=summary,
            tags=[],
            model="local-llm"
        )
        
    def generate_summary_stream(self, content: str) -> Iterator[str]:
        """Generates a summary stream for the given text content."""
        try:
            messages = [
                SystemMessage(content="You are a helpful assistant that summarizes documents efficiently."),
                HumanMessage(content=f"Please provide a concise summary of the following document:\n\n{content[:8000]}")
            ]
            for chunk in self.llm.stream(messages):
                yield chunk.content
        except Exception as e:
            yield f"Error generating summary: {str(e)}"

    def process_file_stream(self, path: str):
        """Reads file, streams summary generation, and stores in DB."""
        logger.info(f"Processing file stream for: {path}")
        if not self.local_file_service or not self.db_manager:
            logger.error("Services not fully initialized for file stream processing")
            yield "Services not fully initialized"
            return

        # 1. Check if summary exists first? 
        # For now, we assume user clicked "Generate" so they want a fresh one or it doesn't exist.
        # But if we want to be smart, we could check DB. 
        # For this task, let's regenerate to show the streaming effect.

        # 2. Get Content
        file_data = self.local_file_service.get_content(path)
        if "error" in file_data:
            logger.warning(f"Error getting file content for {path}: {file_data['error']}")
            yield f"Error: {file_data['error']}"
            return
        
        if file_data.get("type") != "text":
             logger.warning(f"Cannot summarize non-text file: {path}")
             yield "Error: Only text files can be summarized currently."
             return
            
        content = file_data.get("content")
        if not content:
             logger.warning(f"File is empty: {path}")
             yield "Error: File is empty."
             return

        # 3. Stream Summary & Collect
        logger.debug("Starting summary stream")
        full_summary = ""
        for chunk in self.generate_summary_stream(content):
            full_summary += chunk
            yield chunk
        
        # 4. Save to DB after streaming completes
        logger.info(f"Saving summary to database for: {path}")
        self.db_manager.save_summary(
            path=path,
            summary=full_summary,
            tags=[],
            model="local-llm"
        )
        logger.info(f"Summary saved successfully for: {path}")
        
    def generate_tags(self, content: str) -> List[str]:
        """Generates tag suggestions for the given text content."""
        logger.debug(f"Generating tags for content of length: {len(content)}")
        try:
            # Fetch existing tags
            existing_tags = []
            if self.db_manager:
                existing_tags = self.db_manager.get_all_tags()
            
            existing_tags_str = ", ".join(existing_tags) if existing_tags else "None"
            logger.debug(f"Existing tags: {existing_tags_str}")

            messages = [
                SystemMessage(content=f"You are a helpful assistant that analyzes documents and suggests meaningful tags. You have access to the following existing tags: [{existing_tags_str}]. Prioritize using these tags if they are relevant, but you may create new ones if necessary. Provide only a comma-separated list of 3-5 tags. Do not include any other text."),
                HumanMessage(content=f"Please suggest tags for the following document:\n\n{content[:5000]}")
            ]
            response = self.llm.invoke(messages)
            # Parse response
            raw_tags = response.content.strip()
            # Split by comma and clean
            tags = [t.strip() for t in raw_tags.split(',') if t.strip()]
            logger.info(f"Generated {len(tags)} tags: {tags[:5]}")
            return tags[:5] # Limit to 5
        except Exception as e:
            logger.error(f"Error generating tags: {str(e)}", exc_info=True)
            return []

    def process_file_tags(self, path: str):
        """Generates tags for a file."""
        logger.info(f"Processing file tags for: {path}")
        if not self.local_file_service:
             logger.error("LocalFileService not initialized")
             return {"error": "Services not fully initialized"}

        # 1. Get Content
        file_data = self.local_file_service.get_content(path)
        if "error" in file_data:
            logger.warning(f"Error getting file content for tags: {file_data['error']}")
            return file_data
        
        if file_data.get("type") != "text":
             logger.warning(f"Cannot generate tags for non-text file: {path}")
             return {"error": "Only text files can be processed."}
            
        content = file_data.get("content")
        if not content:
             logger.warning(f"File is empty: {path}")
             return {"error": "File is empty."}

        # 2. Generate Tags
        tags = self.generate_tags(content)
        logger.info(f"Tags generated for {path}: {tags}")
        return {"tags": tags}

    def llm_query_with_context(self, query: str, context_text: str):
        messages = [
            SystemMessage(content=(
                "You are a helpful assistant specialized in answering questions based solely on the provided text.\n"
                "1. Answer the user's question **strictly** using only the information inside the <context> tags below.\n"
                "2. If the answer is not contained in the context, explicitly say you don't know. Do not make up an answer.\n"
                "3. Cite the source for your answer using the '(Source: ...)' format provided in the text.\n"
                "4. Do not use emojis. Keep the answer simple, concise, and professional.\n\n"
                f"<context>\n{context_text}\n</context>"
            )),
            HumanMessage(content=query)
        ]
        response = self.llm.invoke(messages)
        return {
            "status": "success",
            "response": response.content.strip()
        }
   

import backend.utils as utils

class RagService:
    def __init__(self, base_url: str = "http://host.docker.internal:1234/v1", embed_llm: str = "text-embedding-granite-embedding-278m-multilingual", debug: bool = True, inmemory: bool = False, root_dir: str = None):
        logger.info(f"Initializing RagService with base_url: {base_url}, embed_llm: {embed_llm}")
        self.db_manager = None
        # Use provided root_dir or fall back to environment variable or default
        if root_dir is None:
            root_dir = os.getenv("DATA_DIR", "./data")
        self.local_file_service = LocalFileService(root_dir=root_dir)
        self.base_url = base_url # TODO: Use environment variables
        self.embed_llm = embed_llm # TODO: Use environment variables or keep as user choice?
        self.embedder = OpenAIEmbeddings(
            base_url=self.base_url,
            api_key="lm-studio", # TODO: Use environment variables
            model=self.embed_llm,
            check_embedding_ctx_length=False
        )
        self.llm = LLMService(base_url=self.base_url)
        if inmemory:
            self.vectorstore = InMemoryVectorStore(self.embedder)
            logger.info("Using in-memory vector store")
        else:
            chroma_base_dir = os.getenv("CHROMA_DIR", "/app/data/chroma")
            self.vectorstore = Chroma(
                persist_directory=chroma_base_dir, 
                embedding_function=self.embedder,
                collection_name="summa_collection"
            )
            logger.info("Using ChromaDB vector store")
        logger.info("RagService initialized successfully")

    def ingest_files(self, paths: List[str]):
        logger.info(f"Ingesting {len(paths)} files into RAG system")
        try:
            utils.test_api()
            docs = []
            for path in paths:
                logger.debug(f"Ingesting file: {path}")
                content = utils.get_file_content(path)
                doc = Document(
                    page_content=content,
                    metadata={"source": path}
                    # TODO: Add chunk indexing for window retrieval
                )
                docs.append(doc)
            chunks = self._split_documents(docs)
            logger.debug(f"Split documents into {len(chunks)} chunks")
            document_ids = self.vectorstore.add_documents(chunks)
            logger.info(f"Successfully ingested {len(paths)} files, {len(chunks)} chunks, {len(document_ids)} document IDs")
            return {
                "status": "success",
                "document_ids": document_ids,
                "message": f"Successfully ingested {len(paths)} files."
            }
        except Exception as e:
            logger.error(f"Error ingesting files: {str(e)}", exc_info=True)
            raise e
    
    def ingest_uploaded_file(self, file: UploadFile):
        logger.info(f"Ingesting uploaded file: {file.filename}")
        try:
            content = self.local_file_service.get_uploaded_file_content(file)
            logger.debug(f"Content: {content}")
            logger.debug(f"Content: {content['content']}")
            logger.debug(f"Metadata: {file.filename}")
            logger.debug(f"Type of content: {type(content)}")
            logger.debug(f"Type of metadata: {type(file.filename)}")
            doc = Document(
                page_content=content['content'],
                metadata={"source": file.filename}
            )
            chunks = self._split_documents([doc])
            logger.debug(f"Split document into {len(chunks)} chunks")
            document_ids = self.vectorstore.add_documents(chunks)
            logger.info(f"Successfully ingested uploaded file, {len(chunks)} chunks, {len(document_ids)} document IDs")
            return {
                "status": "success",
                "document_ids": document_ids,
                "message": f"Successfully ingested uploaded file.",
                "content": content['content'],
                "filename": file.filename
            }
        except Exception as e:
            logger.error(f"Error ingesting uploaded file: {str(e)}", exc_info=True)
            raise e
    
    def _split_documents(self, documents: List[Document]):
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            is_separator_regex=False,
            add_start_index=True
        )
        return splitter.split_documents(documents)
    
    def _vector_search(self, query: str, k: int = 4):
        logger.debug(f"Performing vector search for query: {query[:100]}..., k={k}")
        results = self.vectorstore.similarity_search(query, k=k)
        logger.debug(f"Vector search returned {len(results)} results")
        return results
    
    def query_with_context(self, query: str, k: int = 4):
        logger.info(f"RAG query with context: {query[:100]}...")
        results = self._vector_search(query, k=k)
        for doc in results:
            doc.page_content = f"{doc.page_content} (Source: {doc.metadata['source']})"
        context_text = "\n\n".join([doc.page_content for doc in results])
        logger.debug(f"Context text length: {len(context_text)} chars")
        response = self.llm.llm_query_with_context(query, context_text)
        logger.info("RAG query completed successfully")
        return response



        
